#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CN=RGBColor(0x00,0x3F,0x72); CB=RGBColor(0x00,0x7B,0xC2); CL=RGBColor(0x00,0x9E,0xD9)
CW=RGBColor(0xFF,0xFF,0xFF); COW=RGBColor(0xF4,0xF7,0xFA); CD=RGBColor(0x1C,0x1C,0x2E)
CG=RGBColor(0x6B,0x7A,0x90); CDK=RGBColor(0x13,0x13,0x25); CLG=RGBColor(0xE8,0xED,0xF3)
W=Cm(33.87); H=Cm(19.05); LM=Cm(1.0); CWF=Cm(31.87); HDR=Cm(2.6); FTR=Cm(0.7)
LOGO="/Users/johannesnguyen/Documents/doku/corpdesign/DKFZ_Logo-horiz_de_White_RGB_Avance.png"

prs=Presentation("/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx")

def clear(sl):
    tree=sl.shapes._spTree
    for sh in list(sl.shapes): tree.remove(sh._element)
def rect(sl,l,t,w,h,fill,border=None):
    sp=sl.shapes.add_shape(1,l,t,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if border: sp.line.color.rgb=border; sp.line.width=Pt(0.5)
    else: sp.line.fill.background()
    sp.shadow.inherit=False; return sp
def txt(sl,l,t,w,h,text,sz=14,bold=False,color=CD,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor; tf.margin_top=Cm(0.02); tf.margin_bottom=Cm(0.02)
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Arial"
    return tb
def pic(sl,path,l,t,w):
    try: sl.shapes.add_picture(path,l,t,width=w)
    except Exception as e: print("img:",e)
def card(sl,l,t,w,h,title,lines,tbg=CN,fsz=14):
    th=Cm(0.72); rect(sl,l,t,w,th,tbg)
    txt(sl,l+Cm(0.3),t,w-Cm(0.4),th,title,sz=13,bold=True,color=CW,anchor=MSO_ANCHOR.MIDDLE)
    rect(sl,l,t+th,w,h-th,COW,border=CL)
    n=max(len(lines),1); step=(h-th-Cm(0.3))/n; y=t+th+Cm(0.15)
    for ln in lines:
        txt(sl,l+Cm(0.35),y,w-Cm(0.6),step,f"▸  {ln}",sz=fsz,color=CD,anchor=MSO_ANCHOR.MIDDLE)
        y+=step

s=prs.slides[3]
clear(s)
# chrome
rect(s,0,0,W,HDR,CN); rect(s,0,HDR,W,Cm(0.05),CL)
txt(s,LM,Cm(0.2),W-Cm(7.5),Cm(1.45),"Unternehmen und Projektumfeld",sz=23,bold=True,color=CW)
txt(s,LM,Cm(1.6),W-Cm(7.5),Cm(0.85),"02",sz=12,color=CL)
pic(s,LOGO,W-Cm(6.8),Cm(0.3),Cm(6.2))
rect(s,0,H-FTR,W,FTR,CN)  # Footer einheitlich navy
txt(s,LM,H-FTR+Cm(0.1),W-Cm(4),FTR-Cm(0.15),
    "Johannes Nguyen  |  Fachinformatiker Systemintegration  |  DKFZ / ODCF  |  IHK Rhein-Neckar",sz=8,color=CW)
txt(s,W-Cm(3.5),H-FTR+Cm(0.1),Cm(3.2),FTR-Cm(0.15),"4 / 15",sz=8,color=CW,align=PP_ALIGN.RIGHT)

# two top cards
cw2=(CWF-Cm(0.5))/2
card(s,LM,Cm(3.1),cw2,Cm(6.4),"DKFZ: Deutsches Krebsforschungszentrum",
     ["Größte biomedizinische Forschungseinrichtung Deutschlands",
      "Standort: Heidelberg",
      "Forschung zu Krebsentstehung und -bekämpfung",
      "Komplexe IT-Infrastruktur für genomische Daten"],fsz=14)
card(s,LM+cw2+Cm(0.5),Cm(3.1),cw2,Cm(6.4),"ODCF: Omics IT and Data Management Core Facility",
     ["Interner IT-Dienstleister des DKFZ",
      "Hochleistungs-Compute-Cluster (IBM LSF)",
      "PVE-Cluster: 6 Hypervisor-Knoten, ~60 VMs + LXC",
      "Auftraggeber: Frank Thommen (ODCF)"],tbg=CB,fsz=14)

# Projektrahmen block (bigger)
by=Cm(10.2); bh=Cm(5.6); th=Cm(0.8)
rect(s,LM,by,CWF,th,CN)
txt(s,LM+Cm(0.3),by,CWF-Cm(0.4),th,"Projektrahmen",sz=15,bold=True,color=CW,anchor=MSO_ANCHOR.MIDDLE)
rect(s,LM,by+th,CWF,bh-th,COW,border=CL)
rows=[("Projektart","Internes Projekt, Auftraggeber = Auftragnehmer (DKFZ)"),
      ("Ausbildungsberuf","Fachinformatiker Systemintegration"),
      ("Zeitrahmen","40 Stunden, Bearbeitungszeitraum bis 30.04.2026")]
inner_top=by+th; inner_h=bh-th
pitch=inner_h/len(rows)
y=inner_top
for label,val in rows:
    txt(s,LM+Cm(0.6),y,Cm(7.0),pitch,label,sz=18,bold=True,color=CN,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,LM+Cm(8.0),y,CWF-Cm(8.4),pitch,val,sz=18,color=CD,anchor=MSO_ANCHOR.MIDDLE)
    y+=pitch

prs.save("/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx")
print("saved")
