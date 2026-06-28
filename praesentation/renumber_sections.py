#!/usr/bin/env python3
"""Setzt die Abschnittsnummern der Inhaltsfolien (3..13) auf 01..11,
passend zur Gliederung."""
from pptx import Presentation

EMU = 360000
PPTX = "/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx"
prs = Presentation(PPTX)

for idx in range(2, 13):           # Folien 3..13 (0-basiert 2..12)
    num = f"{idx-1:02d}"           # -> 01..11
    sl = prs.slides[idx]
    target = None
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        t = (sh.top or 0)/EMU
        l = (sh.left or 0)/EMU
        if 1.4 < t < 1.9 and l < 2.0:   # Untertitel-/Nummern-Box unter dem Titel
            target = sh
            break
    if target is None:
        print(f"  Folie {idx+1}: keine Untertitel-Box gefunden")
        continue
    p = target.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = num
        for r in p.runs[1:]:
            r.text = ""
    print(f"  Folie {idx+1} -> {num}")

prs.save(PPTX)
print("saved")
