#!/usr/bin/env python3
"""Rebuild slides 7, 8, 12 cleanly + strip en/em dashes deck-wide."""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── colours ──
CN  = RGBColor(0x00,0x3F,0x72)  # navy
CB  = RGBColor(0x00,0x7B,0xC2)  # blue
CL  = RGBColor(0x00,0x9E,0xD9)  # light blue
CW  = RGBColor(0xFF,0xFF,0xFF)
COW = RGBColor(0xF4,0xF7,0xFA)  # off-white
CROW= RGBColor(0xDD,0xE7,0xF2)  # row stripe (clearly visible)
CD  = RGBColor(0x1C,0x1C,0x2E)  # dark text
CG  = RGBColor(0x6B,0x7A,0x90)  # gray
CGR = RGBColor(0x00,0x9E,0x73)  # green
CR  = RGBColor(0xCC,0x33,0x33)  # red
CAM = RGBColor(0xE6,0x8A,0x00)  # amber
CT  = RGBColor(0x00,0x8F,0xA0)  # teal
CDK = RGBColor(0x13,0x13,0x25)  # footer dark

W   = Cm(33.87)
H   = Cm(19.05)
LM  = Cm(1.0)
CWF = Cm(31.87)        # content width
HDR = Cm(2.6)
FTR = Cm(0.7)
LOGO = "/Users/johannesnguyen/Documents/doku/corpdesign/DKFZ_Logo-horiz_de_White_RGB_Avance.png"

prs = Presentation("/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx")

# ── primitives ──
def clear(slide):
    tree = slide.shapes._spTree
    for sh in list(slide.shapes):
        tree.remove(sh._element)

def rect(sl,l,t,w,h,fill,border=None):
    sp = sl.shapes.add_shape(1,l,t,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border:
        sp.line.color.rgb = border; sp.line.width = Pt(0.5)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def txt(sl,l,t,w,h,text,sz=14,bold=False,color=CD,
        align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Cm(0.02); tf.margin_bottom = Cm(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Arial"
    return tb

def pic(sl,path,l,t,w):
    try: sl.shapes.add_picture(path,l,t,width=w)
    except Exception as e: print("img:",e)

def chrome(sl,title,sub,page):
    rect(sl,0,0,W,HDR,CN)
    rect(sl,0,HDR,W,Cm(0.05),CL)
    txt(sl,LM,Cm(0.2),W-Cm(7.5),Cm(1.45),title,sz=23,bold=True,color=CW)
    if sub:
        txt(sl,LM,Cm(1.6),W-Cm(7.5),Cm(0.85),sub,sz=12,color=CL)
    pic(sl,LOGO,W-Cm(6.8),Cm(0.3),Cm(6.2))
    rect(sl,0,H-FTR,W,FTR,CN)  # Footer einheitlich navy
    txt(sl,LM,H-FTR+Cm(0.1),W-Cm(4),FTR-Cm(0.15),
        "Johannes Nguyen  |  Fachinformatiker Systemintegration  |  DKFZ / ODCF  |  IHK Rhein-Neckar",
        sz=8,color=CW)
    txt(sl,W-Cm(3.5),H-FTR+Cm(0.1),Cm(3.2),FTR-Cm(0.15),
        f"{page} / 15",sz=8,color=CW,align=PP_ALIGN.RIGHT)

def card(sl,l,t,w,h,title,lines,tbg=CN,bbg=COW,tfg=CW,bfg=CD,
         blt="▸",fsz=14,tsz=13):
    th = Cm(0.72)
    rect(sl,l,t,w,th,tbg)
    txt(sl,l+Cm(0.3),t,w-Cm(0.4),th,title,sz=tsz,bold=True,color=tfg,
        anchor=MSO_ANCHOR.MIDDLE)
    rect(sl,l,t+th,w,h-th,bbg,border=CL)
    n = max(len(lines),1)
    avail = h-th-Cm(0.3)
    step = avail/n
    y = t+th+Cm(0.15)
    for ln in lines:
        txt(sl,l+Cm(0.35),y,w-Cm(0.6),step,f"{blt}  {ln}",sz=fsz,color=bfg,
            anchor=MSO_ANCHOR.MIDDLE)
        y += step

def table(sl, x, y, col_w, col_x, header, rows, row_h=Cm(0.82),
          total_w=None, hdr_bg=CN, stripe=True, last_dark=False):
    """Draw a clean table with full-width row backgrounds + cell text."""
    if total_w is None:
        total_w = sum(col_w)+Cm(0.0)
    # header
    rect(sl, x, y, total_w, row_h, hdr_bg)
    for c,(cw,cx) in zip(header, zip(col_w,col_x)):
        txt(sl, cx+Cm(0.2), y, cw-Cm(0.3), row_h, c, sz=12, bold=True,
            color=CW, anchor=MSO_ANCHOR.MIDDLE)
    yy = y+row_h
    for i,row in enumerate(rows):
        is_last = last_dark and i==len(rows)-1
        if is_last:
            bg = CDK; fg = CW; bold=True
        else:
            bg = CROW if (stripe and i%2==0) else CW
            fg = CD; bold=False
        rect(sl, x, yy, total_w, row_h, bg)
        for cell,(cw,cx) in zip(row, zip(col_w,col_x)):
            # allow per-cell color override via tuple (text,color)
            ccol = fg
            ctext = cell
            if isinstance(cell, tuple):
                ctext, ccol = cell
            txt(sl, cx+Cm(0.2), yy, cw-Cm(0.3), row_h, ctext, sz=12,
                bold=bold, color=ccol, anchor=MSO_ANCHOR.MIDDLE)
        yy += row_h
    return yy

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7  – Analysephase
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides[6]
clear(s)
chrome(s, "Analysephase: Anforderungen und Wirtschaftlichkeit", "04", 7)

cw2 = (CWF - Cm(0.5))/2
card(s, LM, Cm(3.1), cw2, Cm(6.0),
     "Funktionale Anforderungen",
     ["Dedizierter PBS auf eigener Hardware",
      "Sicherung aller ~60 VMs und Container",
      "Backup-Jobs mit Zeitplan und Retention",
      "Benachrichtigung bei Fehlschlägen",
      "Anbindung an Prometheus-Monitoring",
      "File- und vollständiger VM-Restore"], fsz=14)
card(s, LM+cw2+Cm(0.5), Cm(3.1), cw2, Cm(6.0),
     "Nicht-funktionale Anforderungen",
     ["Trennung Primär- und Sicherungs-Storage",
      "LDAP-Authentifizierung (DKFZ AD)",
      "Open-Source, keine Lizenzkosten",
      "Integration in bestehendes VLAN",
      "Notfallzugang via lokalem root"], tbg=CB, fsz=14)

# Wirtschaftlichkeit heading
txt(s, LM, Cm(9.45), CWF, Cm(0.65), "Wirtschaftlichkeitsbetrachtung",
    sz=15, bold=True, color=CN)

# cost table (left ~60%)
col_w = [Cm(8.5), Cm(6.5), Cm(4.0)]
col_x = [Cm(1.0), Cm(9.5), Cm(16.0)]
table(s, Cm(1.0), Cm(10.2), col_w, col_x,
      ["Position", "Basis", "Kosten"],
      [["Personalkosten Prüfling", "40 h x ~25 EUR/h", "~1.000 EUR"],
       ["Betreuung", "5 h x ~80 EUR/h", "~400 EUR"],
       ["Hardware", "IBM x3755 (Bestand)", "0 EUR"],
       ["Software", "PBS (AGPLv3)", "0 EUR"],
       ["Betriebskosten", "Strom / Stellplatz", "~200 EUR"],
       ["Gesamt", "", "~1.600 EUR"]],
      row_h=Cm(0.82), total_w=Cm(19.0), last_dark=True)

# Nutzen card (right) – aligned to table top, compact bullets
nx = Cm(20.5); nw = W-nx-Cm(1.0)
card(s, nx, Cm(10.2), nw, Cm(0.82)*7,
     "Nutzen",
     ["Schnellere Restores (Minuten statt Stunden)",
      "Zentrale Verwaltung aller Backups",
      "Automatische Integritätsprüfung",
      "Keine Lizenzkosten"],
     tbg=CGR, fsz=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8  – Begründete Entscheidungen  (2x2, minimal text)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides[7]
clear(s)
chrome(s, "Begründete Entscheidungen", "Warum diese Lösung?", 8)

DEC = [
    (CN,  "PBS statt Veeam / Bacula", "Beauftragt",
     ["Vom Auftraggeber vorgegeben",
      "Keine Lizenzkosten (AGPLv3)",
      "Nativ in PVE integriert"]),
    (CB,  "IBM x3755 M3", "Bestandshardware",
     ["Altbestand, keine Beschaffung",
      "32 Kerne / 32 GB RAM",
      "~20 TB Datastore (RAID 5)"]),
    (CT,  "Hardware-RAID 5 statt ZFS", "Hardware-Vorgabe",
     ["Server kann nur Hardware-RAID",
      "ZFS braucht direkten Disk-Zugriff",
      "ZFS bevorzugt bei neuer Hardware"]),
    (CGR, "ext4 als Dateisystem", "Folgt aus RAID",
     ["ZFS technisch nicht möglich",
      "Vom PBS offiziell unterstützt",
      "Integrität via PBS-Verify-Job"]),
]
gw = (CWF-Cm(0.5))/2
gh = (Cm(15.0)-Cm(0.3))/2
for i,(col,title,tag,bullets) in enumerate(DEC):
    gx = LM + (i%2)*(gw+Cm(0.5))
    gy = Cm(3.0) + (i//2)*(gh+Cm(0.2))
    # title bar
    rect(s, gx, gy, gw, Cm(0.72), col)
    txt(s, gx+Cm(0.3), gy, gw-Cm(0.4), Cm(0.72), title,
        sz=15, bold=True, color=CW, anchor=MSO_ANCHOR.MIDDLE)
    # body
    rect(s, gx, gy+Cm(0.72), gw, gh-Cm(0.72), COW, border=CL)
    # tag chip
    rect(s, gx+Cm(0.4), gy+Cm(0.95), Cm(5.2), Cm(0.6), col)
    txt(s, gx+Cm(0.5), gy+Cm(0.95), Cm(5.0), Cm(0.6), tag,
        sz=11, bold=True, color=CW, anchor=MSO_ANCHOR.MIDDLE)
    # bullets
    by = gy+Cm(1.85)
    step = (gh-Cm(2.1))/len(bullets)
    for b in bullets:
        rect(s, gx+Cm(0.45), by+step/2-Cm(0.16), Cm(0.32), Cm(0.32), col)
        txt(s, gx+Cm(1.0), by, gw-Cm(1.3), step, b, sz=14, color=CD,
            anchor=MSO_ANCHOR.MIDDLE)
        by += step

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Fazit / Soll-Ist-Vergleich
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides[11]
clear(s)
chrome(s, "Fazit: Soll-Ist-Vergleich", "10", 12)

# ECHTE Zahlen aus Tabellen/Zeitnachher.tex (Betriebs- + Projektdoku zusammengefasst):
# Gesamt 40 h geplant / 34,5 h tatsaechlich / -5,5 h. Diff>0 rot, <0 gruen, 0 dunkel.
PD = [("Analyse","4 h","1,5 h","-2,5 h"),
      ("Planung & Konzeption","5 h","3 h","-2 h"),
      ("Beschaffung & Vorbereitung","3 h","5 h","+2 h"),
      ("Installation PBS","4 h","5 h","+1 h"),
      ("Konfiguration Datastore","4 h","2 h","-2 h"),
      ("PVE-Integration","4 h","1 h","-3 h"),
      ("Backup-Jobs","5 h","2 h","-3 h"),
      ("Test & Validierung","3 h","3 h","0 h"),
      ("Monitoring","2 h","2 h","0 h"),
      ("Dokumentation","6 h","10 h","+4 h")]
rows=[]
for ph,g,t,d in PD:
    dc = CR if d.startswith("+") else (CGR if d.startswith("-") else CD)
    rows.append([ph, g, t, (d, dc)])
rows.append([("Gesamt",CW),("40 h",CW),("34,5 h",CW),("-5,5 h",CW)])

col_w=[Cm(13.0),Cm(5.0),Cm(6.0),Cm(6.0)]
col_x=[Cm(1.3), Cm(14.3),Cm(19.3),Cm(25.3)]
end_y = table(s, Cm(1.0), Cm(3.1), col_w, col_x,
      ["Phase","Geplant","Tatsächlich","Differenz"],
      rows, row_h=Cm(0.9), total_w=Cm(30.0), last_dark=True)

# Lessons Learned strip below table
ly = end_y + Cm(0.5)
lh = H - FTR - Cm(0.4) - ly
rect(s, LM, ly, CWF, Cm(0.7), CB)
txt(s, LM+Cm(0.3), ly, Cm(12), Cm(0.7), "Lessons Learned",
    sz=13, bold=True, color=CW, anchor=MSO_ANCHOR.MIDDLE)
rect(s, LM, ly+Cm(0.7), CWF, lh-Cm(0.7), COW, border=CL)
LL = ["Frühe Netzwerk-Abstimmung ist kritisch",
      "Altbestands-Hardware schränkt Dateisystemwahl ein",
      "Betriebsdoku so wertvoll wie die Implementierung",
      "PBS + PVE sehr gut toolunterstützt"]
colw = CWF/4
for i,item in enumerate(LL):
    ix = LM + i*colw
    rect(s, ix+Cm(0.3), ly+Cm(0.95), Cm(0.32), Cm(0.32), CB)
    txt(s, ix+Cm(0.8), ly+Cm(0.8), colw-Cm(1.0), lh-Cm(1.0), item,
        sz=12, color=CD, anchor=MSO_ANCHOR.TOP)

# ════════════════════════════════════════════════════════════════════════════
# DECK-WIDE: strip en/em dashes
# ════════════════════════════════════════════════════════════════════════════
def fix_dashes(t):
    t = t.replace(" – ", ": ").replace(" — ", ": ")
    t = t.replace("–","-").replace("—","-")
    return t

cnt=0
for slide in prs.slides:
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "–" in r.text or "—" in r.text:
                    r.text = fix_dashes(r.text); cnt+=1
print(f"dash fixes: {cnt}")

prs.save("/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx")
print("saved")
