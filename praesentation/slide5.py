#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CN=RGBColor(0x00,0x3F,0x72); CL=RGBColor(0x00,0x9E,0xD9)
CW=RGBColor(0xFF,0xFF,0xFF); CD=RGBColor(0x1C,0x1C,0x2E)
CG=RGBColor(0x6B,0x7A,0x90); CGR=RGBColor(0x00,0x9E,0x73)
CR=RGBColor(0xCC,0x33,0x33); CDK=RGBColor(0x13,0x13,0x25)
CLG=RGBColor(0xE8,0xED,0xF3)
W=Cm(33.87); H=Cm(19.05); LM=Cm(1.0); CWF=Cm(31.87); HDR=Cm(2.6); FTR=Cm(0.7)
LOGO="/Users/johannesnguyen/Documents/doku/corpdesign/DKFZ_Logo-horiz_de_White_RGB_Avance.png"

prs=Presentation("/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx")

def clear(sl):
    tree=sl.shapes._spTree
    for sh in list(sl.shapes): tree.remove(sh._element)
def rect(sl,l,t,w,h,fill,border=None):
    sp=sl.shapes.add_shape(1,l,t,w,h); sp.fill.solid()
    sp.fill.fore_color.rgb=fill
    if border: sp.line.color.rgb=border; sp.line.width=Pt(0.5)
    else: sp.line.fill.background()
    sp.shadow.inherit=False; return sp
def txt(sl,l,t,w,h,text,sz=14,bold=False,color=CD,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor; tf.margin_top=Cm(0.02); tf.margin_bottom=Cm(0.02)
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Arial"
    return tb
def pic(sl,path,l,t,w):
    try: sl.shapes.add_picture(path,l,t,width=w)
    except Exception as e: print("img:",e)

s=prs.slides[4]
clear(s)
# chrome
rect(s,0,0,W,HDR,CN); rect(s,0,HDR,W,Cm(0.05),CL)
txt(s,LM,Cm(0.2),W-Cm(7.5),Cm(1.45),"Ausgangssituation und Projektziel",sz=23,bold=True,color=CW)
txt(s,LM,Cm(1.6),W-Cm(7.5),Cm(0.85),"03",sz=12,color=CL)
pic(s,LOGO,W-Cm(6.8),Cm(0.3),Cm(6.2))
rect(s,0,H-FTR,W,FTR,CDK)
txt(s,LM,H-FTR+Cm(0.1),W-Cm(4),FTR-Cm(0.15),
    "Johannes Nguyen  |  Fachinformatiker Systemintegration  |  DKFZ / ODCF  |  IHK Rhein-Neckar",sz=8,color=CW)
txt(s,W-Cm(3.5),H-FTR+Cm(0.1),Cm(3.2),FTR-Cm(0.15),"5 / 15",sz=8,color=CW,align=PP_ALIGN.RIGHT)

COL=(CWF-Cm(3.0))/2
IST_X=LM
SOLL_X=LM+COL+Cm(3.0)
ARR_X=LM+COL+Cm(0.4)

IST=["Backup nur auf lokalem Hypervisor-Storage",
     "Dell-VxRail-Knoten: kein lokaler Storage",
     "Primär- und Sicherungsdaten gemeinsam abgelegt",
     "Keine zentrale Verwaltung und Überwachung",
     "Keine automatische Integritätsprüfung",
     "Restore aus TSM/Tape: aufwändig und langsam"]
SOLL=["Dedizierter PBS auf eigener Hardware",
      "Anbindung an gesamten PVE-Cluster (6 Knoten)",
      "Inkrementelle Backups mit Deduplizierung",
      "Retention Policies (kritisch / Test)",
      "Zentrales Monitoring via Prometheus",
      "E-Mail-Benachrichtigung bei Fehlern",
      "Betriebsdokumentation für das ODCF-Team"]

def column(x,color,header,items):
    # header (square marker, consistent with bullets)
    rect(s,x,Cm(3.25),Cm(0.5),Cm(0.5),color)
    txt(s,x+Cm(0.7),Cm(3.1),COL-Cm(0.7),Cm(0.8),header,sz=16,bold=True,
        color=color,anchor=MSO_ANCHOR.MIDDLE)
    top=Cm(4.3); bot=Cm(14.4)
    pitch=(bot-top)/len(items)
    y=top
    for it in items:
        rect(s,x+Cm(0.1),y+pitch/2-Cm(0.18),Cm(0.4),Cm(0.4),color)
        txt(s,x+Cm(0.75),y,COL-Cm(0.9),pitch,it,sz=15,color=CD,
            anchor=MSO_ANCHOR.MIDDLE)
        y+=pitch

column(IST_X,CR,"IST-Zustand: Probleme",IST)
column(SOLL_X,CGR,"Projektziel",SOLL)

# arrow centered vertically in items area
txt(s,ARR_X,Cm(8.3),Cm(2.2),Cm(1.5),"→",sz=44,bold=True,color=CL,align=PP_ALIGN.CENTER)

# scope note
rect(s,LM,Cm(14.9),CWF,Cm(0.03),CLG)
txt(s,LM,Cm(15.1),CWF,Cm(0.8),
    "Scope: Blueprint-Testinstanz als Referenz für den späteren Produktivgang, kein Ersatz der bisherigen TSM/Tape-Sicherung",
    sz=12,color=CG)

prs.save("/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx")
print("saved")
