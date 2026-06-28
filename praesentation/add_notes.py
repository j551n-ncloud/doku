#!/usr/bin/env python3
"""Fuegt Sprechernotizen zu allen 15 Folien hinzu (Notizen-Ansicht in PowerPoint)."""
from pptx import Presentation

PPTX = "/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx"
prs = Presentation(PPTX)

notes = {
1: """Begrüßung des Prüfungsausschusses. Kurz vorstellen: Johannes Nguyen, Auszubildender zum Fachinformatiker Systemintegration am DKFZ.
Thema: Einrichtung eines dedizierten Proxmox Backup Servers bei der ODCF.
In den nächsten rund 15 Minuten zeige ich Ausgangslage, Lösung und Ergebnis.""",
2: """Kurzer Überblick über den Ablauf. Nicht vorlesen, nur grob durchgehen.
Erst der technische und betriebliche Kontext, dann Planung und Analyse, anschließend Entwurf und Umsetzung, am Ende Abnahme, Fazit und Ausblick.""",
3: """Damit alle die Begriffe einordnen können, kurz die vier Kerntechnologien.
Proxmox VE: Virtualisierungsplattform, auf der die VMs und Container laufen, bei uns 6 Knoten.
Proxmox Backup Server: die eigentliche Backup-Lösung, inkrementell und mit Deduplizierung.
Ceph: das verteilte Primär-Storage des Clusters.
Hardware-RAID gegen ZFS: wichtig für eine spätere Entscheidung, der alte Server kann nur Hardware-RAID.""",
4: """Das DKFZ ist die größte biomedizinische Forschungseinrichtung Deutschlands in Heidelberg.
Die ODCF ist der interne IT-Dienstleister und betreibt unter anderem den Proxmox-Cluster mit rund 60 VMs.
Wichtig: internes Projekt, Auftraggeber und Auftragnehmer sind beide das DKFZ. Auftraggeber ist Frank Thommen.""",
5: """Die Ausgangslage links: Backups lagen nur lokal auf den Hypervisoren, die VxRail-Knoten hatten gar keinen lokalen Storage. Primär- und Sicherungsdaten lagen zusammen, keine zentrale Verwaltung, keine Integritätsprüfung, Restores über Tape waren langsam.
Das Ziel rechts: ein dedizierter, zentraler Backup-Server mit inkrementellen Backups, Monitoring und Dokumentation.
Hinweis auf den Scope: es ist eine Blueprint-Instanz, kein Ersatz der bestehenden Sicherung.""",
6: """40 Stunden Gesamtbudget, davon maximal 8 Stunden Dokumentation.
Der Balken zeigt die Aufteilung auf die Phasen von Analyse bis Dokumentation.
Vorgehen war iterativ, nach jeder Phase kurze Abstimmung mit dem Betreuer.
Hardware war Bestandshardware, Software komplett Open Source, daher keine Lizenzkosten.""",
7: """Aus Antrag und Ist-Analyse ergeben sich funktionale und nicht-funktionale Anforderungen, links die wichtigsten.
Bei der Wirtschaftlichkeit: rund 1.600 Euro Projektkosten, im Wesentlichen Arbeitszeit, da Hardware und Software nichts kosten.
Der Nutzen rechts: deutlich schnellere Restores, zentrale Verwaltung, Integritätsprüfung. Ein einziger vermiedener langer Restore amortisiert das Projekt.""",
8: """Hier die vier zentralen Entscheidungen, jeweils mit Begründung.
PBS statt Veeam oder Bacula: vom Auftraggeber vorgegeben, keine Lizenzkosten, nativ in PVE integriert.
IBM x3755: Bestandshardware, keine Beschaffung nötig.
Hardware-RAID 5 statt ZFS: der Server kann nur Hardware-RAID, ZFS bräuchte direkten Disk-Zugriff. Bei neuer Hardware wäre ZFS die bevorzugte Lösung.
ext4: folgt zwingend aus der RAID-Entscheidung, Integrität wird über den PBS-Verify-Job sichergestellt.""",
9: """Die Architektur: der PBS hängt im server-VLAN am PVE-Cluster, angebunden über zwei gebündelte 10-Gigabit-Leitungen.
Authentifizierung über das DKFZ-LDAP plus lokaler root für den Notfall.
Monitoring über pbs-exporter und Prometheus, Benachrichtigung per E-Mail.
Backup-Konzept unten: zwei Klassen, beide zweimal täglich, die letzten drei Sicherungen werden aufbewahrt.""",
10: """Die Umsetzung in vier Schritten.
Erst Hardware vorbereiten: RAID-Controller und Netzwerkkarte, hier gab es Mehraufwand durch den RAID-Controller.
Dann PBS installieren und grundkonfigurieren.
Datastore anlegen und LDAP anbinden.
Zuletzt Integration in den PVE-Cluster und Einrichtung der Backup-Jobs, vorher mit einem Test-Backup geprüft.""",
11: """Alle Tests gegen das Backup-Konzept wurden erfolgreich durchlaufen: Backup-Job, LDAP-Login, E-Mail, Monitoring, Storage-Trennung.
Besonders wichtig die Restore-Tests: einzelne Dateien, ganze VM und Container, alle erfolgreich.
Integritätsprüfung über den Verify-Job bestanden.
Abnahme wurde von Frank Thommen erteilt, Restores deutlich schneller als über Tape.""",
12: """Soll-Ist-Vergleich: insgesamt 34,5 statt 40 Stunden, also leicht unter Budget.
Mehraufwand bei Beschaffung und Installation, dafür schneller bei Konfiguration und Anbindung, weil PBS und PVE gut toolunterstützt sind.
Lessons Learned unten: frühe Netzwerk-Abstimmung wichtig, Altbestands-Hardware begrenzt die Dateisystemwahl, gute Betriebsdoku ist genauso wertvoll wie die Umsetzung.""",
13: """Ausblick in drei Horizonten.
Kurzfristig: produktive Instanz auf Basis des Blueprints und schrittweise Migration.
Mittelfristig: Off-Site-Replikation und Ablösung der Tape-Sicherung.
Langfristig: vollautomatisiertes Monitoring und regelmäßige Restore-Tests als Standard.""",
14: """Kurze Zusammenfassung der Kernpunkte: Ziel erreicht, alle Tests bestanden, im Zeitrahmen geblieben, Dokumentation vollständig, klarer Mehrwert für die ODCF.
Überleitung: damit ist die Basis für den späteren Produktivgang gelegt.""",
15: """Dank für die Aufmerksamkeit und Überleitung zum Fachgespräch.
Anbieten, auf Fragen einzugehen, zum Beispiel zu Backup-Konzept, RAID gegen ZFS oder zur LDAP-Anbindung.""",
}

for i, slide in enumerate(prs.slides, 1):
    slide.notes_slide.notes_text_frame.text = notes.get(i, "")

prs.save(PPTX)
print("Notizen zu", len(prs.slides), "Folien hinzugefügt.")
