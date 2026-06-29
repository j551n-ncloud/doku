#!/usr/bin/env python3
"""Baut Folie 12 (Fazit: Soll-Ist-Vergleich) mit den ECHTEN Zahlen aus
Tabellen/Zeitnachher.tex neu: 40 h geplant, 34,5 h tatsaechlich, -5,5 h.
Loest die Diskrepanz zu Folie 14 (die bereits 34,5 / 40 / -5,5 nennt)."""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CN=RGBColor(0x00,0x3F,0x72); CB=RGBColor(0x00,0x7B,0xC2); CL=RGBColor(0x00,0x9E,0xD9)
CW=RGBColor(0xFF,0xFF,0xFF); COW=RGBColor(0xF4,0xF7,0xFA); CROW=RGBColor(0xDD,0xE7,0xF2)
CD=RGBColor(0x1C,0x1C,0x2E); CG=RGBColor(0x6B,0x7A,0x90)
CGR=RGBColor(0x00,0x9E,0x73); CR=RGBColor(0xCC,0x33,0x33)
W=Cm(33.87); H=Cm(19.05); LM=Cm(1.0); CWF=Cm(31.87); HDR=Cm(2.6); FTR=Cm(0.7)
LOGO="/Users/johannesnguyen/Documents/doku/corpdesign/DKFZ_Logo-horiz_de_White_RGB_Avance.png"
PPTX="/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx"

prs=Presentation(PPTX)

def clear(sl):
    tree=sl.shapes._spTree
    for sh in list(sl.shapes): tree.remove(sh._element)
def rect(sl,l,t,w,h,fill,border=None):
    sp=sl.shapes.add_shape(1,l,t,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if border: sp.line.color.rgb=border; sp.line.width=Pt(0.5)
    else: sp.line.fill.background()
    sp.shadow.inherit=False; return sp
def txt(sl,l,t,w,h,text,sz=14,bold=False,color=CD,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor; tf.margin_top=Cm(0.02); tf.margin_bottom=Cm(0.02)
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Arial"
    return tb
def pic(sl,path,l,t,w):
    try: sl.shapes.add_picture(path,l,t,width=w)
    except Exception as e: print("img:",e)

s=prs.slides[11]
clear(s)
# chrome
rect(s,0,0,W,HDR,CN); rect(s,0,HDR,W,Cm(0.05),CL)
txt(s,LM,Cm(0.2),W-Cm(7.5),Cm(1.45),"Fazit: Soll-Ist-Vergleich",sz=23,bold=True,color=CW)
txt(s,LM,Cm(1.6),W-Cm(7.5),Cm(0.85),"10",sz=12,color=CL)
pic(s,LOGO,W-Cm(6.8),Cm(0.3),Cm(6.2))
rect(s,0,H-FTR,W,FTR,CN)
txt(s,LM,H-FTR+Cm(0.1),W-Cm(4),FTR-Cm(0.15),
    "Johannes Nguyen  |  Fachinformatiker Systemintegration  |  DKFZ / ODCF  |  IHK Rhein-Neckar",sz=8,color=CW)
txt(s,W-Cm(3.5),H-FTR+Cm(0.1),Cm(3.2),FTR-Cm(0.15),"12 / 15",sz=8,color=CW,align=PP_ALIGN.RIGHT)

# ECHTE Zahlen aus Zeitnachher.tex (Betriebs- und Projektdoku zu "Dokumentation" zusammengefasst)
# (phase, geplant, tatsaechlich, differenz)  diff>0 rot, diff<0 gruen, 0 dunkel
PD=[("Analyse","4 h","1,5 h","-2,5 h"),
    ("Planung & Konzeption","5 h","3 h","-2 h"),
    ("Beschaffung & Vorbereitung","3 h","5 h","+2 h"),
    ("Installation PBS","4 h","5 h","+1 h"),
    ("Konfiguration Datastore","4 h","2 h","-2 h"),
    ("PVE-Integration","4 h","1 h","-3 h"),
    ("Backup-Jobs","5 h","2 h","-3 h"),
    ("Test & Validierung","3 h","3 h","0 h"),
    ("Monitoring","2 h","2 h","0 h"),
    ("Dokumentation","6 h","10 h","+4 h")]
GESAMT=("Gesamt","40 h","34,5 h","-5,5 h")

col_w=[Cm(13.0),Cm(5.0),Cm(6.0),Cm(6.0)]
col_x=[Cm(1.3),Cm(14.3),Cm(19.3),Cm(25.3)]
total_w=Cm(30.0)
row_h=Cm(0.78)

def diffcolor(d):
    if d.startswith("+"): return CR
    if d.startswith("-"): return CGR
    return CD

y=Cm(3.0)
# header
rect(s,Cm(1.0),y,total_w,row_h,CN)
for c,cw,cx in zip(["Phase","Geplant","Tatsächlich","Differenz"],col_w,col_x):
    txt(s,cx+Cm(0.2),y,cw-Cm(0.3),row_h,c,sz=12,bold=True,color=CW,anchor=MSO_ANCHOR.MIDDLE)
y+=row_h
for i,(ph,g,t,d) in enumerate(PD):
    rect(s,Cm(1.0),y,total_w,row_h,CROW if i%2==0 else CW)
    cells=[(ph,CD),(g,CD),(t,CD),(d,diffcolor(d))]
    for (val,col),cw,cx in zip(cells,col_w,col_x):
        txt(s,cx+Cm(0.2),y,cw-Cm(0.3),row_h,val,sz=12,color=col,anchor=MSO_ANCHOR.MIDDLE)
    y+=row_h
# GESAMT
rect(s,Cm(1.0),y,total_w,row_h,CD)
for (val),cw,cx in zip(GESAMT,col_w,col_x):
    txt(s,cx+Cm(0.2),y,cw-Cm(0.3),row_h,val,sz=12,bold=True,color=CW,anchor=MSO_ANCHOR.MIDDLE)
y+=row_h

# Lessons Learned strip
ly=y+Cm(0.5)
lh=H-FTR-Cm(0.4)-ly
rect(s,LM,ly,CWF,Cm(0.7),CB)
txt(s,LM+Cm(0.3),ly,Cm(12),Cm(0.7),"Lessons Learned",sz=13,bold=True,color=CW,anchor=MSO_ANCHOR.MIDDLE)
rect(s,LM,ly+Cm(0.7),CWF,lh-Cm(0.7),COW,border=CL)
LL=["Frühe Netzwerk-Abstimmung ist kritisch",
    "Altbestands-Hardware schränkt Dateisystemwahl ein",
    "Betriebsdoku so wertvoll wie die Implementierung",
    "PBS + PVE sehr gut toolunterstützt"]
colw=CWF/4
for i,item in enumerate(LL):
    ix=LM+i*colw
    rect(s,ix+Cm(0.3),ly+Cm(0.95),Cm(0.32),Cm(0.32),CB)
    txt(s,ix+Cm(0.8),ly+Cm(0.8),colw-Cm(1.0),lh-Cm(1.0),item,sz=12,color=CD)

prs.save(PPTX)
print("Folie 12 mit echten Zahlen neu gebaut (40 / 34,5 / -5,5).")
