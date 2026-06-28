#!/usr/bin/env python3
"""Vereinheitlicht die Fusszeile aller Folien: Balken navy, Text weiss."""
from pptx import Presentation
from pptx.dml.color import RGBColor

CN = RGBColor(0x00, 0x3F, 0x72)
CW = RGBColor(0xFF, 0xFF, 0xFF)
EMU = 360000
PPTX = "/Users/johannesnguyen/Documents/doku/Praesentaion_Nguyen_Johannes.pptx"

prs = Presentation(PPTX)
bar = txt = 0
for sl in prs.slides:
    for sh in sl.shapes:
        top = (sh.top or 0)/EMU
        left = (sh.left or 0)/EMU
        w = (sh.width or 0)/EMU
        if top > 17.5:
            if sh.has_text_frame and sh.text_frame.text.strip():
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = CW
                txt += 1
            elif left < 0.2 and w > 33.0:
                sh.fill.solid(); sh.fill.fore_color.rgb = CN
                sh.line.fill.background()
                bar += 1
prs.save(PPTX)
print(f"footer bars: {bar}, footer text boxes: {txt}")
